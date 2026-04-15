from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = ROOT / 'deliverables' / 'apresentacao-cliente'
SLIDES_DIR = OUTPUT_DIR / 'slides'
ASSETS_DIR = OUTPUT_DIR / 'assets'
LOGO_PATH = ROOT / 'public' / 'logo.png'
PPTX_PATH = OUTPUT_DIR / 'SDR-Imobiliario-Apresentacao-Premium-Clientes.pptx'

WIDTH = 1920
HEIGHT = 1080

BG_TOP = (9, 13, 24)
BG_BOTTOM = (20, 31, 54)
SURFACE = (18, 27, 43, 220)
SURFACE_ALT = (28, 39, 63, 235)
SURFACE_SOFT = (255, 255, 255, 16)
BORDER = (83, 106, 150, 110)
WHITE = (244, 248, 255, 255)
TEXT = (224, 233, 249, 255)
TEXT_SOFT = (155, 173, 205, 255)
TEXT_DIM = (116, 136, 170, 255)
BRAND = (92, 120, 255, 255)
BRAND_SOFT = (92, 120, 255, 44)
BRAND_ALT = (56, 224, 255, 255)
GREEN = (56, 214, 149, 255)
ORANGE = (255, 179, 71, 255)
RED = (255, 107, 129, 255)
PURPLE = (154, 114, 255, 255)
SHADOW = (2, 6, 14, 150)
WA_HEADER = (12, 129, 97, 255)
WA_BUBBLE_OUT = (214, 255, 206, 255)
WA_BUBBLE_IN = (255, 255, 255, 255)
WA_BG = (231, 241, 235, 255)
WA_TEXT_DARK = (25, 36, 38, 255)
WA_TIME = (101, 118, 123, 255)

FONT_REGULAR = '/System/Library/Fonts/SFNS.ttf'
FONT_BOLD = '/System/Library/Fonts/Supplemental/Arial Bold.ttf'
FONT_HEADLINE = '/System/Library/Fonts/Avenir.ttc'


def rgba(color: tuple[int, int, int] | tuple[int, int, int, int], alpha: int | None = None) -> tuple[int, int, int, int]:
    if len(color) == 4:
        r, g, b, a = color
        return (r, g, b, alpha if alpha is not None else a)
    r, g, b = color
    return (r, g, b, 255 if alpha is None else alpha)


def font(size: int, *, bold: bool = False, headline: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_HEADLINE if headline else FONT_BOLD if bold else FONT_REGULAR
    return ImageFont.truetype(path, size=size)


def create_canvas() -> Image.Image:
    img = Image.new('RGBA', (WIDTH, HEIGHT), BG_TOP)
    draw = ImageDraw.Draw(img, 'RGBA')

    for y in range(HEIGHT):
        t = y / (HEIGHT - 1)
        color = tuple(int(BG_TOP[i] * (1 - t) + BG_BOTTOM[i] * t) for i in range(3))
        draw.line((0, y, WIDTH, y), fill=color)

    glow = Image.new('RGBA', (WIDTH, HEIGHT), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow, 'RGBA')
    gdraw.ellipse((1100, -120, 1820, 600), fill=(70, 120, 255, 80))
    gdraw.ellipse((-180, 640, 520, 1280), fill=(42, 204, 255, 58))
    gdraw.ellipse((620, 180, 1220, 760), fill=(111, 66, 255, 42))
    glow = glow.filter(ImageFilter.GaussianBlur(80))
    img.alpha_composite(glow)

    grid = Image.new('RGBA', (WIDTH, HEIGHT), (0, 0, 0, 0))
    g = ImageDraw.Draw(grid, 'RGBA')
    step = 80
    for x in range(0, WIDTH, step):
        g.line((x, 0, x, HEIGHT), fill=(255, 255, 255, 10))
    for y in range(0, HEIGHT, step):
        g.line((0, y, WIDTH, y), fill=(255, 255, 255, 8))
    grid = grid.filter(ImageFilter.GaussianBlur(0.6))
    img.alpha_composite(grid)
    return img


def add_shadow(base: Image.Image, box: tuple[int, int, int, int], radius: int, *, color: tuple[int, int, int, int] = SHADOW, blur: int = 28, offset: tuple[int, int] = (0, 18)) -> None:
    layer = Image.new('RGBA', base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer, 'RGBA')
    x1, y1, x2, y2 = box
    ox, oy = offset
    draw.rounded_rectangle((x1 + ox, y1 + oy, x2 + ox, y2 + oy), radius=radius, fill=color)
    layer = layer.filter(ImageFilter.GaussianBlur(blur))
    base.alpha_composite(layer)


def draw_card(base: Image.Image, box: tuple[int, int, int, int], *, fill: tuple[int, int, int, int] = SURFACE, outline: tuple[int, int, int, int] = BORDER, radius: int = 28, shadow: bool = True) -> None:
    if shadow:
        add_shadow(base, box, radius)
    draw = ImageDraw.Draw(base, 'RGBA')
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=2)


def rounded_panel(size: tuple[int, int], *, fill: tuple[int, int, int, int] = SURFACE, outline: tuple[int, int, int, int] = BORDER, radius: int = 28) -> Image.Image:
    panel = Image.new('RGBA', size, (0, 0, 0, 0))
    d = ImageDraw.Draw(panel, 'RGBA')
    d.rounded_rectangle((0, 0, size[0] - 1, size[1] - 1), radius=radius, fill=fill, outline=outline, width=2)
    return panel


def paste_logo(base: Image.Image, xy: tuple[int, int], size: int) -> None:
    if not LOGO_PATH.exists():
        return
    logo = Image.open(LOGO_PATH).convert('RGBA')
    logo = ImageOps.contain(logo, (size, size))
    base.alpha_composite(logo, dest=xy)


def measure(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return ['']
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        test = f'{current} {word}'
        if draw.textbbox((0, 0), test, font=fnt)[2] <= max_width:
            current = test
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def draw_wrapped_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, fnt: ImageFont.FreeTypeFont, fill: tuple[int, int, int, int], max_width: int, *, line_gap: int = 10) -> int:
    x, y = xy
    lines = wrap_text(draw, text, fnt, max_width)
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += fnt.size + line_gap
    return y


def draw_kicker(draw: ImageDraw.ImageDraw, text: str, xy: tuple[int, int]) -> None:
    x, y = xy
    fnt = font(20, bold=True)
    w, h = measure(draw, text, fnt)
    draw.rounded_rectangle((x, y, x + w + 28, y + h + 18), radius=26, fill=BRAND_SOFT, outline=rgba(BRAND, 90), width=2)
    draw.text((x + 14, y + 8), text, font=fnt, fill=rgba(BRAND_ALT))


def draw_bullet_group(base: Image.Image, title: str, items: Iterable[str], xy: tuple[int, int], width: int) -> None:
    draw = ImageDraw.Draw(base, 'RGBA')
    x, y = xy
    draw.text((x, y), title, font=font(30, bold=True), fill=WHITE)
    y += 54
    for item in items:
        draw.rounded_rectangle((x, y + 8, x + 18, y + 26), radius=9, fill=rgba(BRAND_ALT, 235))
        draw.text((x + 36, y), item, font=font(22), fill=TEXT_SOFT)
        y = draw_wrapped_text(draw, (x + 36, y), item, font(22), TEXT_SOFT, width - 36, line_gap=6) + 20


def draw_stat_chip(base: Image.Image, label: str, value: str, xy: tuple[int, int], width: int) -> None:
    x, y = xy
    box = (x, y, x + width, y + 94)
    draw_card(base, box, fill=rgba((20, 30, 49), 220), outline=rgba(BRAND_ALT, 38), radius=26, shadow=False)
    draw = ImageDraw.Draw(base, 'RGBA')
    draw.text((x + 22, y + 18), label.upper(), font=font(17, bold=True), fill=TEXT_DIM)
    draw.text((x + 22, y + 44), value, font=font(26, bold=True, headline=True), fill=WHITE)


def draw_topbar(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], title: str) -> None:
    x1, y1, x2, _ = box
    draw.text((x1 + 28, y1 + 22), title, font=font(32, bold=True), fill=WHITE)
    draw.text((x1 + 28, y1 + 64), 'Corretor Digital 24/7 para a operação comercial', font=font(16), fill=TEXT_DIM)
    draw.rounded_rectangle((x2 - 260, y1 + 18, x2 - 110, y1 + 56), radius=18, fill=rgba((255, 255, 255), 18), outline=rgba((255, 255, 255), 26), width=1)
    draw.text((x2 - 238, y1 + 29), 'Buscar', font=font(14), fill=TEXT_DIM)
    draw.ellipse((x2 - 74, y1 + 16, x2 - 34, y1 + 56), fill=rgba(BRAND, 140))


def build_sidebar(height: int, active: str) -> Image.Image:
    sidebar = Image.new('RGBA', (244, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(sidebar, 'RGBA')
    draw.rounded_rectangle((0, 0, 243, height), radius=0, fill=(14, 21, 35, 240))
    paste_logo(sidebar, (28, 22), 48)
    draw.text((92, 24), 'SDR Imobiliário', font=font(22, bold=True), fill=WHITE)
    draw.text((92, 55), 'Corretor Digital 24/7', font=font(13), fill=TEXT_DIM)

    items = ['Visão Geral', 'Mensagens', 'Leads', 'Catálogo', 'WhatsApp', 'Importação']
    y = 128
    for item in items:
        selected = item == active
        if selected:
            draw.rounded_rectangle((18, y - 6, 224, y + 40), radius=18, fill=BRAND_SOFT, outline=rgba(BRAND, 85), width=2)
        fill = WHITE if selected else TEXT_SOFT
        icon_fill = rgba(BRAND_ALT if selected else (180, 192, 218), 255)
        draw.ellipse((32, y + 4, 56, y + 28), fill=icon_fill)
        draw.text((72, y), item, font=font(18, bold=selected), fill=fill)
        y += 62

    draw.rounded_rectangle((20, height - 138, 224, height - 28), radius=22, fill=rgba((255, 255, 255), 14), outline=rgba((255, 255, 255), 22), width=1)
    draw.text((38, height - 120), 'Operação conectada', font=font(15, bold=True), fill=WHITE)
    draw.text((38, height - 90), 'WhatsApp, leads e catálogo', font=font(14), fill=TEXT_DIM)
    draw.rounded_rectangle((38, height - 56, 142, height - 28), radius=14, fill=rgba(GREEN, 34), outline=rgba(GREEN, 88), width=1)
    draw.text((52, height - 50), 'ONLINE', font=font(12, bold=True), fill=rgba(GREEN))
    return sidebar


def base_screen(width: int, height: int, title: str, active: str) -> tuple[Image.Image, ImageDraw.ImageDraw, tuple[int, int, int, int]]:
    screen = Image.new('RGBA', (width, height), (13, 19, 31, 255))
    draw = ImageDraw.Draw(screen, 'RGBA')
    sidebar = build_sidebar(height, active)
    screen.alpha_composite(sidebar, dest=(0, 0))

    header_box = (244, 0, width, 110)
    draw.rounded_rectangle((244, 0, width, 110), radius=0, fill=rgba((255, 255, 255), 8))
    draw_topbar(draw, header_box, title)

    content_box = (272, 136, width - 28, height - 28)
    return screen, draw, content_box


def draw_dashboard_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Visão Geral', 'Visão Geral')
    x1, y1, x2, _ = content

    labels = [
        ('Mensagens Hoje', '148', BRAND_ALT),
        ('Leads Captados', '72', PURPLE),
        ('Leads Quentes', '18', ORANGE),
        ('WhatsApp', 'Conectado', GREEN),
    ]

    card_w = 255
    gap = 18
    for idx, (label, value, color) in enumerate(labels):
        x = x1 + idx * (card_w + gap)
        box = (x, y1, x + card_w, y1 + 166)
        draw_card(screen, box, fill=rgba((23, 34, 56), 245), outline=rgba(color, 90))
        draw.text((x + 22, y1 + 22), label.upper(), font=font(15, bold=True), fill=TEXT_DIM)
        draw.text((x + 22, y1 + 72), value, font=font(48, bold=True, headline=True), fill=WHITE)
        draw.ellipse((x + 190, y1 + 24, x + 228, y1 + 62), fill=rgba(color, 60))
        draw.ellipse((x + 201, y1 + 35, x + 217, y1 + 51), fill=rgba(color))

    actions_y = y1 + 206
    for idx, text in enumerate([
        'Configurar WhatsApp',
        'Ver Mensagens',
        'Gerenciar Leads',
    ]):
        x = x1 + idx * 340
        box = (x, actions_y, x + 312, actions_y + 130)
        draw_card(screen, box, fill=rgba((20, 31, 49), 240), outline=rgba(BRAND, 70), radius=24, shadow=False)
        draw.rounded_rectangle((x + 20, actions_y + 22, x + 64, actions_y + 66), radius=14, fill=BRAND_SOFT, outline=rgba(BRAND, 60), width=1)
        draw.text((x + 82, actions_y + 26), text, font=font(22, bold=True), fill=WHITE)
        draw.text((x + 82, actions_y + 60), 'Ação rápida dentro do painel', font=font(15), fill=TEXT_DIM)

    chart_box = (x1, actions_y + 170, x2 - 380, actions_y + 470)
    draw_card(screen, chart_box, fill=rgba((18, 27, 44), 240), outline=rgba(BRAND_ALT, 42), radius=26, shadow=False)
    draw.text((chart_box[0] + 26, chart_box[1] + 20), 'Pipeline semanal', font=font(22, bold=True), fill=WHITE)
    draw.text((chart_box[0] + 26, chart_box[1] + 54), 'Operação com visão comercial em tempo real', font=font(15), fill=TEXT_DIM)

    chart_left = chart_box[0] + 36
    chart_bottom = chart_box[3] - 34
    bars = [120, 160, 190, 154, 210, 236, 264]
    for idx, h in enumerate(bars):
        x = chart_left + idx * 86
        draw.rounded_rectangle((x, chart_bottom - h, x + 42, chart_bottom), radius=16, fill=rgba(BRAND, 220))
        draw.rounded_rectangle((x + 48, chart_bottom - h + 42, x + 82, chart_bottom), radius=14, fill=rgba(BRAND_ALT, 180))
        draw.text((x, chart_bottom + 8), ['Seg', 'Ter', 'Qua', 'Qui', 'Sex', 'Sáb', 'Dom'][idx], font=font(13), fill=TEXT_DIM)

    right_box = (x2 - 340, actions_y + 170, x2, actions_y + 470)
    draw_card(screen, right_box, fill=rgba((18, 27, 44), 240), outline=rgba(PURPLE, 42), radius=26, shadow=False)
    draw.text((right_box[0] + 24, right_box[1] + 20), 'Resumo comercial', font=font(22, bold=True), fill=WHITE)
    rows = [
        ('Tempo médio de resposta', 'instantâneo'),
        ('Leads qualificados hoje', '12'),
        ('Catálogo disponível', '218 imóveis'),
        ('WhatsApp da operação', 'online'),
    ]
    y = right_box[1] + 76
    for label, value in rows:
        draw.rounded_rectangle((right_box[0] + 24, y, right_box[2] - 24, y + 52), radius=16, fill=rgba((255, 255, 255), 10))
        draw.text((right_box[0] + 42, y + 15), label, font=font(14), fill=TEXT_SOFT)
        vw, _ = measure(draw, value, font(15, bold=True))
        draw.text((right_box[2] - 40 - vw, y + 14), value, font=font(15, bold=True), fill=WHITE)
        y += 66

    return screen


def draw_messages_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Mensagens', 'Mensagens')
    x1, y1, x2, y2 = content

    left_box = (x1, y1, x1 + 332, y2)
    draw_card(screen, left_box, fill=rgba((20, 30, 49), 242), outline=rgba(BRAND, 40), radius=28, shadow=False)
    draw.text((x1 + 26, y1 + 22), 'Conversas', font=font(24, bold=True), fill=WHITE)
    draw.rounded_rectangle((x1 + 24, y1 + 70, x1 + 308, y1 + 114), radius=18, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 18), width=1)
    draw.text((x1 + 48, y1 + 83), 'Buscar...', font=font(15), fill=TEXT_DIM)

    convs = [
        ('Mariana Souza', '(11) 98765-4321', 'Procuro apê 3 quartos', 'agora', 'Quente', True),
        ('Carlos Mendes', '(21) 99812-1100', 'Tem opções na Barra?', '6 min', 'Frio', False),
        ('Ana Ribeiro', '(31) 99122-8899', 'Quero agendar visita', '18 min', 'Agendou Visita', False),
        ('Felipe Lima', '(47) 99900-0022', 'Qual o valor do condomínio?', '40 min', 'Frio', False),
    ]
    row_y = y1 + 138
    for name, phone, preview, ago, status, active in convs:
        if active:
            draw.rounded_rectangle((x1 + 14, row_y - 6, x1 + 318, row_y + 82), radius=20, fill=rgba(BRAND, 28), outline=rgba(BRAND, 90), width=2)
        draw.ellipse((x1 + 28, row_y + 4, x1 + 66, row_y + 42), fill=rgba(BRAND_ALT if active else (180, 192, 218), 70))
        draw.text((x1 + 82, row_y + 2), name, font=font(18, bold=True), fill=WHITE)
        draw.text((x1 + 82, row_y + 28), phone, font=font(13), fill=TEXT_DIM)
        draw.text((x1 + 82, row_y + 50), preview, font=font(13), fill=TEXT_SOFT)
        draw.text((x1 + 260, row_y + 2), ago, font=font(12), fill=TEXT_DIM)
        badge_color = GREEN if status == 'Agendou Visita' else ORANGE if status == 'Quente' else BRAND_ALT
        bw, _ = measure(draw, status, font(11, bold=True))
        draw.rounded_rectangle((x1 + 210, row_y + 46, x1 + 226 + bw, row_y + 72), radius=12, fill=rgba(badge_color, 32), outline=rgba(badge_color, 70), width=1)
        draw.text((x1 + 222, row_y + 52), status, font=font(11, bold=True), fill=rgba(badge_color))
        row_y += 94

    right_box = (x1 + 356, y1, x2, y2)
    draw_card(screen, right_box, fill=rgba((18, 28, 44), 242), outline=rgba(BRAND_ALT, 44), radius=28, shadow=False)
    draw.text((right_box[0] + 26, y1 + 24), 'Mariana Souza', font=font(24, bold=True), fill=WHITE)
    draw.text((right_box[0] + 26, y1 + 56), '(11) 98765-4321', font=font(14), fill=TEXT_DIM)

    badge_text = 'Quente'
    bw, _ = measure(draw, badge_text, font(12, bold=True))
    bx = right_box[0] + 280
    draw.rounded_rectangle((bx, y1 + 22, bx + bw + 30, y1 + 50), radius=13, fill=rgba(ORANGE, 36), outline=rgba(ORANGE, 82), width=1)
    draw.text((bx + 16, y1 + 29), badge_text, font=font(12, bold=True), fill=rgba(ORANGE))

    danger = 'Limpar conversa'
    dw, _ = measure(draw, danger, font(13, bold=True))
    dx = right_box[2] - dw - 74
    draw.rounded_rectangle((dx, y1 + 16, right_box[2] - 20, y1 + 56), radius=16, fill=rgba(RED, 20), outline=rgba(RED, 70), width=1)
    draw.text((dx + 18, y1 + 28), danger, font=font(13, bold=True), fill=rgba(RED))

    bubbles = [
        ('in', 'Oi, queria ver opções de apartamento em Moema com 3 quartos.'),
        ('out', 'Claro. Posso te mostrar opções com vaga e lazer também. Qual faixa de investimento você procura?'),
        ('in', 'Até R$ 1,4 milhão, com pelo menos 2 vagas.'),
        ('out', 'Perfeito. Já separei algumas opções no perfil que você descreveu e posso te encaminhar para um corretor.'),
    ]
    bubble_y = y1 + 124
    for direction, text in bubbles:
        max_width = 420
        bubble_font = font(15)
        lines = wrap_text(draw, text, bubble_font, max_width)
        bubble_h = 28 + len(lines) * 22 + 22
        if direction == 'out':
            x = right_box[2] - max_width - 54
            fill = rgba(BRAND, 220)
            text_fill = WHITE
        else:
            x = right_box[0] + 26
            fill = rgba((255, 255, 255), 12)
            text_fill = TEXT
        bubble_w = max(min(max(measure(draw, line, bubble_font)[0] for line in lines) + 42, max_width + 42), 230)
        draw.rounded_rectangle((x, bubble_y, x + bubble_w, bubble_y + bubble_h), radius=24, fill=fill, outline=rgba((255, 255, 255), 18), width=1)
        ty = bubble_y + 18
        for line in lines:
            draw.text((x + 20, ty), line, font=bubble_font, fill=text_fill)
            ty += 22
        draw.text((x + 20, bubble_y + bubble_h - 24), '09:24', font=font(11), fill=rgba((255, 255, 255), 170 if direction == 'out' else 110))
        bubble_y += bubble_h + 18

    return screen


def draw_leads_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Leads', 'Leads')
    x1, y1, x2, y2 = content

    draw.text((x1, y1), 'Leads qualificados pelo Corretor Digital', font=font(18), fill=TEXT_DIM)
    chips = [('Todos (72)', True), ('Frio (36)', False), ('Quente (24)', False), ('Agendou Visita (12)', False)]
    cx = x1
    for label, active in chips:
        w, _ = measure(draw, label, font(14, bold=True))
        draw.rounded_rectangle((cx, y1 + 34, cx + w + 34, y1 + 72), radius=18, fill=rgba(BRAND, 30) if active else rgba((255, 255, 255), 8), outline=rgba(BRAND, 60) if active else rgba((255, 255, 255), 16), width=1)
        draw.text((cx + 18, y1 + 46), label, font=font(14, bold=True), fill=WHITE if active else TEXT_SOFT)
        cx += w + 48

    draw.rounded_rectangle((x1, y1 + 98, x1 + 368, y1 + 144), radius=18, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 18), width=1)
    draw.text((x1 + 22, y1 + 112), 'Buscar por nome, telefone...', font=font(15), fill=TEXT_DIM)

    table_box = (x1, y1 + 174, x2, y2)
    draw_card(screen, table_box, fill=rgba((18, 28, 44), 242), outline=rgba(BRAND_ALT, 40), radius=28, shadow=False)

    headers = ['Lead', 'Telefone', 'Orçamento', 'Bairro', 'Status', 'Msgs', 'Último contato', 'Ações']
    cols = [290, 180, 170, 160, 130, 70, 170, 150]
    cursor = x1 + 24
    for header, width in zip(headers, cols):
        draw.text((cursor, table_box[1] + 20), header.upper(), font=font(12, bold=True), fill=TEXT_DIM)
        cursor += width
    draw.line((x1 + 24, table_box[1] + 54, x2 - 24, table_box[1] + 54), fill=rgba((255, 255, 255), 18), width=1)

    leads = [
        ('Mariana Souza', '(11) 98765-4321', 'R$ 1.400.000', 'Moema', 'Quente', '14', 'agora'),
        ('Carlos Mendes', '(21) 99812-1100', 'R$ 980.000', 'Barra', 'Frio', '6', '18 min'),
        ('Ana Ribeiro', '(31) 99122-8899', 'R$ 1.850.000', 'Lourdes', 'Agendou Visita', '10', '32 min'),
        ('Felipe Lima', '(47) 99900-0022', 'R$ 760.000', 'Centro', 'Frio', '5', '1 h'),
    ]
    row_y = table_box[1] + 76
    for name, phone, budget, bairro, status, msgs, ago in leads:
        draw.rounded_rectangle((x1 + 18, row_y - 8, x2 - 18, row_y + 52), radius=18, fill=rgba((255, 255, 255), 6))
        cursor = x1 + 24
        draw.ellipse((cursor, row_y - 2, cursor + 28, row_y + 26), fill=rgba(BRAND, 65))
        draw.text((cursor + 42, row_y), name, font=font(16, bold=True), fill=WHITE)
        cursor += cols[0]
        draw.text((cursor, row_y), phone, font=font(14), fill=TEXT_SOFT)
        cursor += cols[1]
        draw.text((cursor, row_y), budget, font=font(14), fill=TEXT_SOFT)
        cursor += cols[2]
        draw.text((cursor, row_y), bairro, font=font(14), fill=TEXT_SOFT)
        cursor += cols[3]
        badge_color = GREEN if status == 'Agendou Visita' else ORANGE if status == 'Quente' else BRAND_ALT
        bw, _ = measure(draw, status, font(12, bold=True))
        draw.rounded_rectangle((cursor, row_y - 2, cursor + bw + 26, row_y + 24), radius=13, fill=rgba(badge_color, 28), outline=rgba(badge_color, 70), width=1)
        draw.text((cursor + 13, row_y + 4), status, font=font(12, bold=True), fill=rgba(badge_color))
        cursor += cols[4]
        draw.text((cursor + 12, row_y), msgs, font=font(14), fill=TEXT_SOFT)
        cursor += cols[5]
        draw.text((cursor, row_y), ago, font=font(14), fill=TEXT_SOFT)
        cursor += cols[6]
        draw.text((cursor, row_y), 'Ver chat', font=font(13, bold=True), fill=rgba(BRAND_ALT))
        draw.text((cursor + 74, row_y), 'Excluir', font=font(13, bold=True), fill=rgba(RED))
        row_y += 74

    return screen


def property_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], title: str, bairro: str, price: str, badges: list[str]) -> None:
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=24, fill=rgba((22, 33, 53), 248), outline=rgba((255, 255, 255), 14), width=1)
    draw.text((x1 + 18, y1 + 16), title, font=font(16, bold=True), fill=WHITE)
    draw.text((x1 + 18, y1 + 40), f'📍 {bairro}', font=font(12), fill=TEXT_DIM)
    draw.text((x1 + 18, y1 + 70), price, font=font(28, bold=True, headline=True), fill=rgba(BRAND_ALT))
    draw.text((x1 + 18, y1 + 114), '3Q   •   2V   •   128m²', font=font(13), fill=TEXT_SOFT)
    cx = x1 + 18
    for badge in badges:
        bw, _ = measure(draw, badge, font(11, bold=True))
        draw.rounded_rectangle((cx, y1 + 144, cx + bw + 18, y1 + 168), radius=12, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 14), width=1)
        draw.text((cx + 9, y1 + 150), badge, font=font(11, bold=True), fill=TEXT_SOFT)
        cx += bw + 28


def draw_catalog_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Catálogo de Imóveis', 'Catálogo')
    x1, y1, x2, _ = content
    draw.text((x1, y1), '218 imóveis no catálogo', font=font(18), fill=TEXT_DIM)
    draw.rounded_rectangle((x1, y1 + 42, x1 + 320, y1 + 86), radius=18, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 16), width=1)
    draw.text((x1 + 20, y1 + 56), 'Buscar por título ou bairro...', font=font(14), fill=TEXT_DIM)
    for i, label in enumerate(['Todos os bairros', 'Quartos', 'Limpar']):
        bx = x1 + 340 + i * 170
        draw.rounded_rectangle((bx, y1 + 42, bx + 150, y1 + 86), radius=18, fill=rgba((255, 255, 255), 9), outline=rgba((255, 255, 255), 16), width=1)
        draw.text((bx + 18, y1 + 56), label, font=font(14), fill=TEXT_SOFT)

    card_w = 350
    card_h = 240
    gap_x = 22
    gap_y = 24
    items = [
        ('Apartamento Alto Padrão', 'Moema', 'R$ 1.250.000', ['Pet OK', 'Piscina']),
        ('Casa com Área Gourmet', 'Alphaville', 'R$ 2.480.000', ['Portaria 24h', 'Lazer']),
        ('Cobertura Duplex', 'Jardins', 'R$ 3.100.000', ['Academia', 'Piscina']),
        ('Apartamento Compacto', 'Centro', 'R$ 680.000', ['Pet OK']),
    ]
    start_y = y1 + 122
    for idx, item in enumerate(items):
        col = idx % 3
        row = idx // 3
        box = (
            x1 + col * (card_w + gap_x),
            start_y + row * (card_h + gap_y),
            x1 + col * (card_w + gap_x) + card_w,
            start_y + row * (card_h + gap_y) + card_h,
        )
        property_card(draw, box, *item)
    return screen


def draw_import_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Importação Administrativa em Lote', 'Importação')
    x1, y1, x2, y2 = content

    draw.text((x1, y1), 'Importe imóveis fornecendo o ID do tenant e um JSON Array', font=font(18), fill=TEXT_DIM)
    main_box = (x1, y1 + 42, x2 - 120, y2 - 80)
    draw_card(screen, main_box, fill=rgba((18, 28, 44), 240), outline=rgba(BRAND, 46), radius=30, shadow=False)

    draw.text((main_box[0] + 28, main_box[1] + 28), 'Selecione o Cliente (Tenant) de Destino', font=font(16, bold=True), fill=WHITE)
    draw.rounded_rectangle((main_box[0] + 28, main_box[1] + 62, main_box[0] + 520, main_box[1] + 108), radius=18, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 18), width=1)
    draw.text((main_box[0] + 48, main_box[1] + 76), 'Selecione um cliente', font=font(14), fill=TEXT_SOFT)

    draw.text((main_box[0] + 28, main_box[1] + 140), 'Lista de Imóveis (JSON Array)', font=font(16, bold=True), fill=WHITE)
    draw.rounded_rectangle((main_box[0] + 28, main_box[1] + 176, main_box[2] - 28, main_box[2] - 28), radius=22, fill=(14, 20, 33, 255), outline=rgba((255, 255, 255), 16), width=1)
    code_y = main_box[1] + 208
    code_lines = [
        '[',
        '  {',
        '    "titulo": "Apartamento Alto Padrão",',
        '    "valor": 1250000,',
        '    "quartos": 3,',
        '    "vagas": 2,',
        '    "link_site": "https://meusite.com/imoveis/1"',
        '  }',
        ']',
    ]
    mono = font(18)
    for line in code_lines:
        draw.text((main_box[0] + 56, code_y), line, font=mono, fill=rgba(BRAND_ALT if '"' in line else TEXT_SOFT))
        code_y += 34

    upload_text = 'Fazer Upload de Arquivo .json'
    uw, _ = measure(draw, upload_text, font(14, bold=True))
    draw.rounded_rectangle((main_box[2] - uw - 84, main_box[1] + 136, main_box[2] - 28, main_box[1] + 176), radius=16, fill=rgba((255, 255, 255), 12), outline=rgba((255, 255, 255), 18), width=1)
    draw.text((main_box[2] - uw - 56, main_box[1] + 149), upload_text, font=font(14, bold=True), fill=TEXT_SOFT)

    action_text = 'Importar Imóveis'
    aw, _ = measure(draw, action_text, font(16, bold=True))
    draw.rounded_rectangle((main_box[0] + 28, main_box[3] - 68, main_box[0] + aw + 90, main_box[3] - 20), radius=18, fill=rgba(BRAND, 226), outline=rgba(BRAND_ALT, 120), width=1)
    draw.text((main_box[0] + 52, main_box[3] - 54), action_text, font=font(16, bold=True), fill=WHITE)
    return screen


def draw_whatsapp_screen() -> Image.Image:
    screen, draw, content = base_screen(1440, 900, 'Conexão WhatsApp', 'WhatsApp')
    x1, y1, x2, y2 = content

    draw.text((x1, y1), 'Conecte o número da imobiliária para ativar o Corretor Digital', font=font(18), fill=TEXT_DIM)

    status_box = (x1, y1 + 46, x2 - 120, y1 + 290)
    draw_card(screen, status_box, fill=rgba((18, 28, 44), 242), outline=rgba(GREEN, 40), radius=30, shadow=False)
    draw.ellipse((status_box[0] + 30, status_box[1] + 34, status_box[0] + 94, status_box[1] + 98), fill=rgba(GREEN, 26), outline=rgba(GREEN, 80), width=2)
    draw.text((status_box[0] + 124, status_box[1] + 36), 'WhatsApp Conectado', font=font(28, bold=True), fill=WHITE)
    draw.text((status_box[0] + 124, status_box[1] + 76), '+55 11 4002-8922', font=font(17), fill=TEXT_SOFT)
    draw.text((status_box[0] + 124, status_box[1] + 104), 'Imobiliária Demo', font=font(14), fill=TEXT_DIM)
    draw.ellipse((status_box[2] - 52, status_box[1] + 52, status_box[2] - 32, status_box[1] + 72), fill=rgba(GREEN))

    draw.rounded_rectangle((status_box[0] + 34, status_box[3] - 70, status_box[0] + 282, status_box[3] - 24), radius=18, fill=rgba(BRAND, 226), outline=rgba(BRAND_ALT, 90), width=1)
    draw.text((status_box[0] + 62, status_box[3] - 56), 'Conectar WhatsApp', font=font(16, bold=True), fill=WHITE)
    draw.rounded_rectangle((status_box[0] + 300, status_box[3] - 70, status_box[0] + 474, status_box[3] - 24), radius=18, fill=rgba((255, 255, 255), 10), outline=rgba((255, 255, 255), 18), width=1)
    draw.text((status_box[0] + 344, status_box[3] - 56), 'Atualizar', font=font(16, bold=True), fill=TEXT_SOFT)

    qr_box = (x1, y1 + 326, x1 + 450, y2 - 16)
    draw_card(screen, qr_box, fill=rgba((18, 28, 44), 240), outline=rgba(BRAND_ALT, 40), radius=28, shadow=False)
    draw.text((qr_box[0] + 28, qr_box[1] + 24), 'Escaneie o QR Code', font=font(24, bold=True), fill=WHITE)
    draw.text((qr_box[0] + 28, qr_box[1] + 58), 'Conecte o dispositivo em segundos', font=font(15), fill=TEXT_DIM)
    draw.rounded_rectangle((qr_box[0] + 84, qr_box[1] + 106, qr_box[0] + 364, qr_box[1] + 386), radius=24, fill=WHITE)
    for i in range(4):
        for j in range(4):
            if (i + j) % 2 == 0:
                sx = qr_box[0] + 114 + i * 54
                sy = qr_box[1] + 136 + j * 54
                draw.rounded_rectangle((sx, sy, sx + 34, sy + 34), radius=6, fill=(12, 16, 28, 255))

    how_box = (x1 + 480, y1 + 326, x2 - 120, y2 - 16)
    draw_card(screen, how_box, fill=rgba((18, 28, 44), 240), outline=rgba(PURPLE, 34), radius=28, shadow=False)
    draw.text((how_box[0] + 28, how_box[1] + 24), 'Como funciona', font=font(24, bold=True), fill=WHITE)
    steps = [
        'Clique em "Conectar WhatsApp" para gerar o QR Code',
        'No celular, abra o WhatsApp e vá em Dispositivos conectados',
        'Escaneie o QR Code com a câmera do app',
        'Pronto. O bot começa a responder automaticamente',
    ]
    sy = how_box[1] + 78
    for idx, step in enumerate(steps):
        draw.ellipse((how_box[0] + 28, sy + 2, how_box[0] + 56, sy + 30), fill=BRAND_SOFT, outline=rgba(BRAND_ALT, 70), width=1)
        draw.text((how_box[0] + 38, sy + 7), str(idx + 1), font=font(13, bold=True), fill=rgba(BRAND_ALT))
        draw_wrapped_text(draw, (how_box[0] + 74, sy), step, font(15), TEXT_SOFT, how_box[2] - how_box[0] - 112, line_gap=8)
        sy += 86
    return screen


def draw_wa_text_bubble(draw: ImageDraw.ImageDraw, width: int, y: int, side: str, text: str, time_text: str) -> int:
    bubble_font = font(17)
    lines = wrap_text(draw, text, bubble_font, 260)
    text_w = max(measure(draw, line, bubble_font)[0] for line in lines)
    bubble_w = max(170, min(292, text_w + 42))
    bubble_h = 22 + len(lines) * 24 + 28
    x = 18 if side == 'left' else width - bubble_w - 18
    fill = WA_BUBBLE_IN if side == 'left' else WA_BUBBLE_OUT

    draw.rounded_rectangle((x, y, x + bubble_w, y + bubble_h), radius=22, fill=fill, outline=(211, 222, 216, 255), width=1)
    ty = y + 14
    for line in lines:
        draw.text((x + 16, ty), line, font=bubble_font, fill=WA_TEXT_DARK)
        ty += 24

    tw, _ = measure(draw, time_text, font(11))
    draw.text((x + bubble_w - tw - 12, y + bubble_h - 18), time_text, font=font(11), fill=WA_TIME)
    return bubble_h


def draw_wa_property_card(draw: ImageDraw.ImageDraw, width: int, y: int, title: str, bairro: str, price: str, meta: str) -> int:
    card_w = 312
    card_h = 186
    x = width - card_w - 18
    draw.rounded_rectangle((x, y, x + card_w, y + card_h), radius=24, fill=WA_BUBBLE_OUT, outline=(189, 217, 190, 255), width=1)
    draw.rounded_rectangle((x + 14, y + 14, x + card_w - 14, y + 74), radius=18, fill=(168, 211, 187, 255))
    draw.text((x + 26, y + 92), title, font=font(16, bold=True), fill=WA_TEXT_DARK)
    draw.text((x + 26, y + 116), bairro, font=font(13), fill=WA_TIME)
    draw.text((x + 26, y + 140), price, font=font(22, bold=True, headline=True), fill=(20, 109, 72, 255))
    draw.text((x + 26, y + 168), meta, font=font(12), fill=WA_TIME)
    return card_h


def draw_wa_system_chip(draw: ImageDraw.ImageDraw, width: int, y: int, text: str) -> int:
    chip_font = font(13, bold=True)
    tw, th = measure(draw, text, chip_font)
    x = (width - tw - 28) // 2
    draw.rounded_rectangle((x, y, x + tw + 28, y + th + 12), radius=14, fill=(219, 228, 223, 255))
    draw.text((x + 14, y + 6), text, font=chip_font, fill=WA_TIME)
    return th + 12


def draw_whatsapp_chat_mockup(contact_name: str, stage_label: str, items: list[dict[str, str]]) -> Image.Image:
    width, height = 460, 960
    img = Image.new('RGBA', (width, height), WA_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    for y in range(height):
        t = y / (height - 1)
        c = tuple(int(WA_BG[i] * (1 - t) + (214, 232, 220)[i] * t) for i in range(3))
        draw.line((0, y, width, y), fill=c)

    pattern = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    pdraw = ImageDraw.Draw(pattern, 'RGBA')
    for x in range(-40, width + 40, 56):
        for y in range(110, height, 56):
            pdraw.ellipse((x, y, x + 12, y + 12), fill=(255, 255, 255, 30))
    pattern = pattern.filter(ImageFilter.GaussianBlur(1))
    img.alpha_composite(pattern)

    draw.rectangle((0, 0, width, 40), fill=(244, 247, 245, 255))
    draw.text((18, 12), '09:41', font=font(14, bold=True), fill=WA_TEXT_DARK)
    draw.text((width - 94, 12), 'LTE', font=font(12, bold=True), fill=WA_TEXT_DARK)
    draw.rounded_rectangle((width - 54, 14, width - 24, 26), radius=4, outline=WA_TEXT_DARK, width=2)

    draw.rectangle((0, 40, width, 116), fill=WA_HEADER)
    draw.text((20, 64), '‹', font=font(26, bold=True), fill=WHITE)
    draw.ellipse((52, 56, 100, 104), fill=(46, 168, 128, 255))
    initials = ''.join(word[0] for word in contact_name.split()[:2]).upper()
    draw.text((64, 69), initials, font=font(18, bold=True), fill=WHITE)
    draw.text((116, 58), contact_name, font=font(19, bold=True), fill=WHITE)
    draw.text((116, 84), stage_label, font=font(12), fill=(220, 246, 236, 255))

    y = 140
    for item in items:
        kind = item.get('kind', 'text')
        if kind == 'system':
            y += draw_wa_system_chip(draw, width, y, item['text']) + 18
            continue
        if kind == 'property':
            y += draw_wa_property_card(draw, width, y, item['title'], item['bairro'], item['price'], item['meta']) + 16
            continue
        y += draw_wa_text_bubble(draw, width, y, item['side'], item['text'], item['time']) + 14

    input_y = height - 86
    draw.rounded_rectangle((16, input_y, width - 16, height - 16), radius=28, fill=(244, 247, 245, 255), outline=(207, 221, 212, 255), width=1)
    draw.text((42, input_y + 21), 'Mensagem', font=font(16), fill=(131, 145, 139, 255))
    draw.ellipse((width - 72, input_y + 14, width - 28, input_y + 58), fill=WA_HEADER)
    draw.text((width - 56, input_y + 22), '›', font=font(24, bold=True), fill=WHITE)
    return img


def paste_device(base: Image.Image, content: Image.Image, box: tuple[int, int, int, int], *, phone: bool = False) -> None:
    x1, y1, x2, y2 = box
    frame = Image.new('RGBA', (x2 - x1, y2 - y1), (0, 0, 0, 0))
    draw = ImageDraw.Draw(frame, 'RGBA')
    radius = 40 if not phone else 54
    draw.rounded_rectangle((0, 0, frame.width - 1, frame.height - 1), radius=radius, fill=(12, 16, 26, 255), outline=rgba((255, 255, 255), 34), width=2)

    inset = 26 if not phone else 22
    inner = ImageOps.fit(content, (frame.width - inset * 2, frame.height - inset * 2), method=Image.Resampling.LANCZOS)
    inner_mask = Image.new('L', inner.size, 0)
    ImageDraw.Draw(inner_mask).rounded_rectangle((0, 0, inner.width, inner.height), radius=28 if not phone else 42, fill=255)
    frame.paste(inner, (inset, inset), inner_mask)

    if not phone:
        draw.ellipse((frame.width // 2 - 6, 10, frame.width // 2 + 6, 22), fill=rgba((255, 255, 255), 90))
        base_shadow_box = (x1, y2 - 12, x2, y2 + 28)
        add_shadow(base, base_shadow_box, radius=20, blur=18, offset=(0, 10))
    else:
        draw.rounded_rectangle((frame.width // 2 - 52, 10, frame.width // 2 + 52, 28), radius=14, fill=rgba((255, 255, 255), 18))

    add_shadow(base, box, radius=radius, blur=34, offset=(0, 20))
    base.alpha_composite(frame, dest=(x1, y1))

    if not phone:
        bezel = Image.new('RGBA', (frame.width + 80, 52), (0, 0, 0, 0))
        bdraw = ImageDraw.Draw(bezel, 'RGBA')
        bdraw.rounded_rectangle((40, 0, bezel.width - 40, 38), radius=18, fill=(174, 182, 203, 90))
        bdraw.rounded_rectangle((0, 16, bezel.width, 52), radius=20, fill=(136, 144, 165, 76))
        bezel = bezel.filter(ImageFilter.GaussianBlur(1))
        base.alpha_composite(bezel, dest=(x1 - 40, y2 - 2))


def add_footer(draw: ImageDraw.ImageDraw, page: int) -> None:
    draw.text((72, HEIGHT - 48), 'Telas ilustrativas baseadas no sistema SDR Imobiliário', font=font(14), fill=TEXT_DIM)
    page_text = f'{page:02d}'
    w, _ = measure(draw, page_text, font(16, bold=True))
    draw.text((WIDTH - 72 - w, HEIGHT - 48), page_text, font=font(16, bold=True), fill=TEXT_DIM)


def slide_cover(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')

    draw_kicker(draw, 'APRESENTAÇÃO COMERCIAL • VERSÃO PREMIUM', (92, 84))
    paste_logo(img, (92, 154), 58)
    draw.text((168, 162), 'SDR Imobiliário', font=font(24, bold=True), fill=WHITE)
    draw.text((168, 197), 'Corretor Digital 24/7 para imobiliárias', font=font(14), fill=TEXT_DIM)

    title = 'Um sistema comercial elegante para captar, atender e converter leads no WhatsApp'
    draw_wrapped_text(draw, (92, 278), title, font(54, bold=True, headline=True), WHITE, 720, line_gap=18)
    subtitle = 'Material pronto para apresentação a clientes, com mockups do painel, da operação comercial e da jornada digital de atendimento.'
    draw_wrapped_text(draw, (92, 486), subtitle, font(24), TEXT_SOFT, 620, line_gap=12)

    draw_stat_chip(img, 'Atendimento', 'em segundos', (92, 640), 222)
    draw_stat_chip(img, 'Leads', 'qualificados', (332, 640), 222)
    draw_stat_chip(img, 'Catálogo', 'centralizado', (572, 640), 222)

    paste_device(img, assets['dashboard'], (960, 140, 1784, 742))
    paste_device(img, assets['wa_intro'], (1334, 442, 1658, 1010), phone=True)
    paste_device(img, assets['wa_visit'], (1590, 514, 1870, 1008), phone=True)

    overlay_box = (1006, 774, 1586, 942)
    draw_card(img, overlay_box, fill=rgba((14, 20, 35), 220), outline=rgba(BRAND_ALT, 40), radius=28)
    draw.text((1038, 804), 'Conversas de WhatsApp como peça central da venda', font=font(22, bold=True), fill=WHITE)
    draw_wrapped_text(draw, (1038, 842), 'Além do painel, a apresentação agora mostra exemplos de atendimento, qualificação e agendamento em mockups de conversa.', font(18), TEXT_SOFT, 500, line_gap=8)

    add_footer(draw, 1)
    return img


def slide_dashboard(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'VISÃO GERAL DA OPERAÇÃO', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Um painel executivo que transforma atividade em clareza comercial', font(48, bold=True, headline=True), WHITE, 720, line_gap=16)
    draw_wrapped_text(draw, (92, 320), 'A primeira impressão para o cliente é de controle: métricas, ações rápidas e leitura visual da operação em tempo real.', font(22), TEXT_SOFT, 660, line_gap=10)
    draw_bullet_group(
        img,
        'Destaques da tela',
        [
            'Indicadores de mensagens, leads e status do WhatsApp.',
            'Blocos de ação rápida para acelerar o time comercial.',
            'Resumo visual que passa sensação de organização e performance.',
        ],
        (92, 432),
        620,
    )
    paste_device(img, assets['dashboard'], (890, 130, 1800, 862))
    add_footer(draw, 2)
    return img


def slide_messages(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'MENSAGENS E ATENDIMENTO 24/7', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Inbox de conversas com leitura rápida e contexto comercial', font(48, bold=True, headline=True), WHITE, 700, line_gap=16)
    draw_wrapped_text(draw, (92, 320), 'O atendimento no WhatsApp ganha cara de operação séria: fila organizada, histórico claro e ação imediata sobre cada conversa.', font(22), TEXT_SOFT, 640, line_gap=10)
    draw_bullet_group(
        img,
        'O que o cliente percebe',
        [
            'Atendimento contínuo com linguagem comercial.',
            'Limpeza de conversa e organização operacional.',
            'Status de lead visível junto do histórico do chat.',
        ],
        (92, 442),
        620,
    )
    paste_device(img, assets['messages'], (822, 118, 1804, 900))
    add_footer(draw, 3)
    return img


def slide_whatsapp_examples(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'MOCKUPS DE CONVERSAS NO WHATSAPP', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'O que mais chama a atenção do cliente: a conversa pronta acontecendo no WhatsApp', font(48, bold=True, headline=True), WHITE, 1040, line_gap=16)
    draw_wrapped_text(draw, (92, 320), 'Em vez de mostrar só dashboard, este slide vende a experiência final. O cliente enxerga abordagem, qualificação, oferta de imóveis e agendamento em linguagem real de atendimento.', font(22), TEXT_SOFT, 1120, line_gap=10)

    captions = [
        ('Resposta imediata', 'Primeiro contato acolhedor, com leitura do perfil do lead e início da qualificação.'),
        ('Sugestão de imóveis', 'Entrega de opções compatíveis com o orçamento, mantendo ritmo comercial.'),
        ('Agendamento e handoff', 'Conversa evolui para visita ou repasse ao corretor no momento certo.'),
    ]
    phones = ['wa_intro', 'wa_options', 'wa_visit']
    x = 106
    for idx, key in enumerate(phones):
        paste_device(img, assets[key], (x, 430, x + 430, 972), phone=True)
        title, desc = captions[idx]
        draw.text((x + 18, 988), title, font=font(20, bold=True), fill=WHITE)
        draw_wrapped_text(draw, (x + 18, 1018), desc, font(15), TEXT_SOFT, 384, line_gap=6)
        x += 560

    add_footer(draw, 4)
    return img


def slide_leads(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'GESTÃO DE LEADS E PRIORIZAÇÃO', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Leads qualificados, organizados e prontos para ação comercial', font(48, bold=True, headline=True), WHITE, 680, line_gap=16)
    draw_wrapped_text(draw, (92, 322), 'A plataforma sai do atendimento e entra na gestão: prioriza quente, organiza agenda e deixa a equipe sempre com próximo passo claro.', font(22), TEXT_SOFT, 630, line_gap=10)

    metrics_box = (92, 450, 660, 816)
    draw_card(img, metrics_box, fill=rgba((17, 26, 42), 224), outline=rgba(PURPLE, 40), radius=28)
    draw.text((124, 484), 'Fluxo comercial organizado', font=font(30, bold=True), fill=WHITE)
    draw.text((124, 532), 'Quente, frio ou visita agendada: o lead já entra no painel com prioridade visual.', font=font(18), fill=TEXT_SOFT)

    y = 610
    for label, value, color in [
        ('Leads em andamento', '72', BRAND_ALT),
        ('Oportunidades quentes', '24', ORANGE),
        ('Visitas agendadas', '12', GREEN),
    ]:
        draw.rounded_rectangle((124, y, 628, y + 62), radius=18, fill=rgba((255, 255, 255), 8), outline=rgba((255, 255, 255), 14), width=1)
        draw.text((148, y + 18), label, font=font(16), fill=TEXT_SOFT)
        vw, _ = measure(draw, value, font(24, bold=True, headline=True))
        draw.text((610 - vw, y + 12), value, font=font(24, bold=True, headline=True), fill=rgba(color))
        y += 82

    paste_device(img, assets['leads'], (760, 116, 1810, 896))
    add_footer(draw, 5)
    return img


def slide_catalog_import(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'CATÁLOGO E IMPORTAÇÃO', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Imóveis organizados em vitrine digital e entrada rápida de novos produtos', font(48, bold=True, headline=True), WHITE, 760, line_gap=16)
    draw_wrapped_text(draw, (92, 324), 'Além do atendimento, a apresentação mostra estrutura: catálogo pesquisável, cards elegantes e importação em lote para escalar operação.', font(22), TEXT_SOFT, 700, line_gap=10)

    left_box = (92, 452, 730, 892)
    draw_card(img, left_box, fill=rgba((17, 26, 42), 224), outline=rgba(BRAND_ALT, 40), radius=28)
    draw.text((124, 486), 'Diferenciais percebidos', font=font(30, bold=True), fill=WHITE)
    draw_bullet_group(
        img,
        '',
        [
            'Catálogo premium com filtros por bairro e dormitórios.',
            'Importação em JSON para entrada rápida de inventário.',
            'Links diretos para o imóvel e apresentação comercial consistente.',
        ],
        (124, 548),
        560,
    )

    paste_device(img, assets['catalog'], (820, 130, 1776, 528))
    paste_device(img, assets['import'], (920, 552, 1710, 952))
    add_footer(draw, 6)
    return img


def slide_whatsapp(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'CONEXÃO COM WHATSAPP', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Onboarding simples para colocar a operação no ar com agilidade', font(48, bold=True, headline=True), WHITE, 700, line_gap=16)
    draw_wrapped_text(draw, (92, 320), 'A conexão do WhatsApp aparece como um passo controlado, com QR Code, status visual e orientação clara para a equipe.', font(22), TEXT_SOFT, 640, line_gap=10)

    draw_bullet_group(
        img,
        'Mensagem que isso passa ao cliente',
        [
            'Implantação rápida e guiada.',
            'Operação conectada com monitoramento visual.',
            'Experiência pronta para escalar atendimento sem perder controle.',
        ],
        (92, 440),
        620,
    )

    paste_device(img, assets['whatsapp'], (820, 120, 1798, 898))
    add_footer(draw, 7)
    return img


def slide_closing(assets: dict[str, Image.Image]) -> Image.Image:
    img = create_canvas()
    draw = ImageDraw.Draw(img, 'RGBA')
    draw_kicker(draw, 'ENCERRAMENTO COMERCIAL', (92, 84))
    draw_wrapped_text(draw, (92, 170), 'Uma apresentação pronta para mostrar velocidade, organização e valor percebido', font(52, bold=True, headline=True), WHITE, 900, line_gap=18)
    draw_wrapped_text(draw, (92, 352), 'Podemos personalizar esta versão com a identidade da imobiliária, catálogo real, argumentos de venda e narrativa comercial específica para cada cliente.', font(24), TEXT_SOFT, 780, line_gap=12)

    steps_box = (92, 542, 980, 900)
    draw_card(img, steps_box, fill=rgba((16, 24, 39), 226), outline=rgba(BRAND_ALT, 42), radius=30)
    draw.text((124, 576), 'Fluxo que a apresentação vende', font=font(32, bold=True), fill=WHITE)
    steps = [
        '1. Conecta o WhatsApp da operação.',
        '2. Centraliza conversas e qualifica o lead.',
        '3. Organiza oportunidades em painel comercial.',
        '4. Mostra catálogo com aparência premium.',
    ]
    sy = 646
    for step in steps:
        draw.text((124, sy), step, font=font(22), fill=TEXT_SOFT)
        sy += 62

    collage = Image.new('RGBA', (760, 420), (0, 0, 0, 0))
    collage.alpha_composite(ImageOps.fit(assets['dashboard'], (440, 260), method=Image.Resampling.LANCZOS), dest=(0, 0))
    collage.alpha_composite(ImageOps.fit(assets['wa_options'], (300, 210), method=Image.Resampling.LANCZOS), dest=(460, 24))
    collage.alpha_composite(ImageOps.fit(assets['leads'], (420, 220), method=Image.Resampling.LANCZOS), dest=(56, 192))
    collage.alpha_composite(ImageOps.fit(assets['catalog'], (280, 176), method=Image.Resampling.LANCZOS), dest=(480, 226))

    collage_box = (1080, 214, 1848, 842)
    draw_card(img, collage_box, fill=rgba((16, 24, 39), 208), outline=rgba(PURPLE, 40), radius=34)
    add_shadow(img, collage_box, radius=34, blur=38, offset=(0, 24))
    img.alpha_composite(collage, dest=(1084, 218))
    add_footer(draw, 8)
    return img


def save_slides(slides: list[Image.Image]) -> list[Path]:
    paths: list[Path] = []
    for idx, slide in enumerate(slides, start=1):
        path = SLIDES_DIR / f'slide-{idx:02d}.png'
        slide.convert('RGB').save(path, quality=95)
        paths.append(path)
    return paths


def build_pptx(slide_paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    for slide_path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(slide_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    if prs.slides:
        first = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(first)

    prs.save(PPTX_PATH)


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    assets = {
        'dashboard': draw_dashboard_screen(),
        'messages': draw_messages_screen(),
        'leads': draw_leads_screen(),
        'catalog': draw_catalog_screen(),
        'import': draw_import_screen(),
        'whatsapp': draw_whatsapp_screen(),
        'wa_intro': draw_whatsapp_chat_mockup(
            'Mariana Souza',
            'online agora',
            [
                {'kind': 'system', 'text': 'Hoje, 09:41'},
                {'side': 'left', 'text': 'Oi, vi vocês no Instagram. Estou procurando apartamento em Moema.', 'time': '09:41'},
                {'side': 'right', 'text': 'Olá, Mariana. Posso te ajudar com isso. Você busca compra ou locação?', 'time': '09:41'},
                {'side': 'left', 'text': 'Compra. Até R$ 1,4 milhão e com 2 vagas.', 'time': '09:42'},
                {'side': 'right', 'text': 'Perfeito. Já estou filtrando opções no seu perfil para te mostrar as melhores alternativas.', 'time': '09:42'},
            ],
        ),
        'wa_options': draw_whatsapp_chat_mockup(
            'Mariana Souza',
            'qualificação concluída',
            [
                {'kind': 'system', 'text': 'Hoje, 09:44'},
                {'side': 'right', 'text': 'Separei duas opções com 3 quartos, 2 vagas e lazer. Veja a primeira:', 'time': '09:44'},
                {'kind': 'property', 'title': 'Apartamento Alto Padrão', 'bairro': 'Moema • 128m²', 'price': 'R$ 1.250.000', 'meta': '3 quartos   •   2 vagas'},
                {'side': 'left', 'text': 'Gostei bastante. Essa tem varanda?', 'time': '09:45'},
                {'side': 'right', 'text': 'Tem sim, varanda gourmet e condomínio com academia e piscina.', 'time': '09:45'},
            ],
        ),
        'wa_visit': draw_whatsapp_chat_mockup(
            'Mariana Souza',
            'lead qualificado',
            [
                {'kind': 'system', 'text': 'Hoje, 09:47'},
                {'side': 'left', 'text': 'Perfeito. Quero agendar visita para sábado.', 'time': '09:47'},
                {'side': 'right', 'text': 'Ótimo. Vou encaminhar agora para um corretor e ele segue com o agendamento.', 'time': '09:47'},
                {'kind': 'system', 'text': 'Lead classificado como Quente'},
                {'side': 'right', 'text': 'Posso confirmar o melhor horário entre manhã e tarde?', 'time': '09:48'},
                {'side': 'left', 'text': 'De manhã, por favor.', 'time': '09:48'},
            ],
        ),
    }

    for name, asset in assets.items():
        asset.convert('RGB').save(ASSETS_DIR / f'{name}.png', quality=95)

    slides = [
        slide_cover(assets),
        slide_dashboard(assets),
        slide_messages(assets),
        slide_whatsapp_examples(assets),
        slide_leads(assets),
        slide_catalog_import(assets),
        slide_whatsapp(assets),
        slide_closing(assets),
    ]
    slide_paths = save_slides(slides)
    build_pptx(slide_paths)
    print(PPTX_PATH)


if __name__ == '__main__':
    main()
